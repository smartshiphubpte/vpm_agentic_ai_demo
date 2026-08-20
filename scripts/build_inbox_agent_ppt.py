#!/usr/bin/env python3
"""Build a 1-slide dummy-friendly Inbox Agent visualization."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "inbox_agent_for_dummies.pptx"

NAVY = RGBColor(0x0B, 0x2E, 0x4A)
TEAL = RGBColor(0x00, 0x7A, 0x8C)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
ORANGE = RGBColor(0xE8, 0x71, 0x0A)
RED = RGBColor(0xC5, 0x22, 0x1F)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
SLATE = RGBColor(0x33, 0x44, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xE8, 0xF4, 0xF8)
LIGHT = RGBColor(0xF0, 0xF6, 0xFA)
ARROW = RGBColor(0x5F, 0x6B, 0x7A)


def _shape(slide, kind, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(kind, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.25)
    else:
        s.line.fill.background()
    return s


def _label(slide, left, top, width, height, text, size=11, bold=False, color=SLATE, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return box


def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    _shape(s, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.82), NAVY)
    _label(s, Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.4),
           "Inbox Agent  —  Email In, Clean Data Out", size=24, bold=True, color=WHITE)
    _label(s, Inches(0.4), Inches(0.48), Inches(12.5), Inches(0.28),
           "Dummy version: a careful mail clerk that never sleeps, never double-files, and never silently throws mail away",
           size=12, color=PALE)

    # One-liner
    _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.95), Inches(12.73), Inches(0.48), TEAL)
    _label(s, Inches(0.4), Inches(0.95), Inches(12.53), Inches(0.48),
           "The job:  grab new mail  →  figure out what it is  →  pull the numbers  →  check them  →  save once",
           size=14, bold=True, color=WHITE)

    # Pipeline steps
    steps = [
        (NAVY, "1  Grab mail", "Look only for NEW emails.\nRemember where we left off."),
        (TEAL, "2  Open it", "Read the body and any\nExcel / CSV / PDF files."),
        (BLUE, "3  Sort it", "Noon? Weather? Port?\nIncident? Unknown?"),
        (ORANGE, "4  Pull data", "Turn messy text into\na clean form."),
        (GREEN, "5  Check it", "Required fields? Right ship?\nUnits and time OK?"),
        (PURPLE, "6  Save once", "Write to the database.\nIf we already have it, skip."),
    ]
    box_w, box_h = Inches(1.82), Inches(2.05)
    gap = Inches(0.22)
    start_x = Inches(0.35)
    y = Inches(1.58)
    for i, (color, title, body) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h, color)
        _label(s, x + Inches(0.06), y + Inches(0.12), box_w - Inches(0.12), Inches(0.5),
               title, size=13, bold=True, color=WHITE)
        _label(s, x + Inches(0.08), y + Inches(0.68), box_w - Inches(0.16), Inches(1.2),
               body, size=11, color=WHITE)
        if i < len(steps) - 1:
            ax = x + box_w
            _shape(s, MSO_SHAPE.RIGHT_ARROW, ax, y + Inches(0.88), gap, Inches(0.28), ARROW)

    # Safety net
    _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(3.8), Inches(12.73), Inches(1.55), LIGHT, TEAL)
    _label(s, Inches(0.45), Inches(3.86), Inches(12.4), Inches(0.32),
           "Safety net  —  if something goes wrong, we do NOT lose the email",
           size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    nets = [
        (ORANGE, "Retry", "Temporary glitch\n(network / DB).\nTry again."),
        (RED, "Dead letter", "Broken / poison mail.\nPark it. Keep it."),
        (BLUE, "Replay", "A person can say\n“try this one again”."),
        (TEAL, "Watch", "Alerts if the pile grows\nor the mailbox dies."),
    ]
    nw, nh = Inches(2.95), Inches(1.05)
    ny = Inches(4.22)
    for i, (color, title, body) in enumerate(nets):
        x = Inches(0.5) + i * (nw + Inches(0.16))
        _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, ny, Inches(0.12), nh, color)
        _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.12), ny, nw - Inches(0.12), nh, WHITE, color)
        _label(s, x + Inches(0.22), ny + Inches(0.04), nw - Inches(0.3), Inches(0.28),
               title, size=13, bold=True, color=color, align=PP_ALIGN.LEFT)
        _label(s, x + Inches(0.22), ny + Inches(0.32), nw - Inches(0.3), Inches(0.68),
               body, size=11, color=SLATE, align=PP_ALIGN.LEFT)

    # Why a week
    _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(5.5), Inches(12.73), Inches(1.78), WHITE, NAVY)
    _label(s, Inches(0.45), Inches(5.55), Inches(12.4), Inches(0.32),
           "Why this takes a week  —  the hard part is NOT “read email”",
           size=14, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    _label(s, Inches(0.45), Inches(5.85), Inches(12.4), Inches(0.24),
           "Safe automation, every day, in production:",
           size=12, color=SLATE, align=PP_ALIGN.LEFT)

    reasons = [
        (TEAL, "Messy mail", "Forwarded chains, HTML,\nweird attachments."),
        (BLUE, "No duplicates", "Same email twice must\nnot create two records."),
        (GREEN, "No silent fail", "If parse breaks, we park\nit and tell someone."),
        (ORANGE, "Proof it worked", "Logs, metrics, alerts,\nreplay for support."),
    ]
    rw, rh = Inches(2.95), Inches(0.95)
    ry = Inches(6.16)
    for i, (color, title, body) in enumerate(reasons):
        x = Inches(0.5) + i * (rw + Inches(0.16))
        _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, ry, rw, rh, color)
        _label(s, x + Inches(0.1), ry + Inches(0.06), rw - Inches(0.2), Inches(0.28),
               title, size=13, bold=True, color=WHITE)
        _label(s, x + Inches(0.1), ry + Inches(0.36), rw - Inches(0.2), Inches(0.52),
               body, size=11, color=WHITE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
