#!/usr/bin/env python3
"""Build pictorial 3-slide pre-voyage weather workflow deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "pre_voyage_weather_workflow.pptx"

NAVY = RGBColor(0x0B, 0x2E, 0x4A)
TEAL = RGBColor(0x00, 0x7A, 0x8C)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
ORANGE = RGBColor(0xE8, 0x71, 0x0A)
RED = RGBColor(0xC5, 0x22, 0x1F)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
SLATE = RGBColor(0x33, 0x44, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF6, 0xFA)
PALE = RGBColor(0xE8, 0xF4, 0xF8)
ARROW = RGBColor(0x5F, 0x6B, 0x7A)


def _shape(slide, kind, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(kind, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def _label(slide, left, top, width, height, text, size=11, bold=False, color=SLATE, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _header(slide, title, subtitle=""):
    _shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.72), NAVY)
    _label(slide, Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.35), title, size=20, bold=True, color=WHITE)
    if subtitle:
        _label(slide, Inches(0.4), Inches(0.42), Inches(9.2), Inches(0.22), subtitle, size=10, color=PALE)


def _agent(slide, left, top, name, role, color):
    w, h = Inches(1.55), Inches(0.95)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, color, NAVY)
    _label(slide, left, top + Inches(0.12), w, Inches(0.35), name, size=9, bold=True, color=WHITE)
    _label(slide, left + Inches(0.05), top + Inches(0.48), w - Inches(0.1), Inches(0.4), role, size=7, color=WHITE)
    return left + w / 2, top + h / 2


def _artifact(slide, left, top, label, ext, color):
    w, h = Inches(1.05), Inches(0.72)
    _shape(slide, MSO_SHAPE.FLOWCHART_DOCUMENT, left, top, w, h, color, SLATE)
    _label(slide, left, top + Inches(0.08), w, Inches(0.3), ext, size=8, bold=True, color=WHITE)
    _label(slide, left, top + Inches(0.38), w, Inches(0.28), label, size=6, color=WHITE)
    return left + w / 2, top + h / 2


def _store(slide, left, top, label):
    w, h = Inches(1.35), Inches(0.9)
    _shape(slide, MSO_SHAPE.FLOWCHART_DATA, left, top, w, h, PURPLE, NAVY)
    _label(slide, left, top + Inches(0.28), w, Inches(0.35), label, size=8, bold=True, color=WHITE)
    return left + w / 2, top + h / 2


def _arrow_h(slide, x1, y, x2, label=""):
    if x2 < x1:
        x1, x2 = x2, x1
    mid_y = y - Inches(0.12)
    _shape(slide, MSO_SHAPE.RIGHT_ARROW, x1, mid_y, x2 - x1, Inches(0.24), ARROW)
    if label:
        _label(slide, (x1 + x2) / 2 - Inches(0.55), y - Inches(0.42), Inches(1.1), Inches(0.22), label, size=7, color=TEAL)


def _arrow_v(slide, x, y1, y2, label=""):
    mid_x = x - Inches(0.1)
    _shape(slide, MSO_SHAPE.DOWN_ARROW, mid_x, y1, Inches(0.2), y2 - y1, ARROW)
    if label:
        _label(slide, x + Inches(0.12), (y1 + y2) / 2 - Inches(0.15), Inches(0.9), Inches(0.3), label, size=7, color=TEAL, align=PP_ALIGN.LEFT)


def _api_cloud(slide, left, top):
    w, h = Inches(1.4), Inches(0.75)
    _shape(slide, MSO_SHAPE.CLOUD, left, top, w, h, BLUE, NAVY)
    _label(slide, left, top + Inches(0.22), w, Inches(0.3), "Weather API", size=8, bold=True, color=WHITE)
    return left + w / 2, top + h / 2


def slide1(prs):
    """End-to-end agent pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Agent Pipeline", "Excel drop → voyage → waypoints → weather → map & report")

    # Row 1: trigger
    _artifact(slide, Inches(0.35), Inches(1.15), "pre_voyage", ".xlsx", GREEN)
    _arrow_h(slide, Inches(1.45), Inches(1.5), Inches(1.85), "file path")
    _agent(slide, Inches(1.85), Inches(1.2), "InboxWatch\nAgent", "polls inbox", TEAL)
    _arrow_h(slide, Inches(3.45), Inches(1.5), Inches(3.85), "dispatch")
    _agent(slide, Inches(3.85), Inches(1.2), "Flow\nRunner", "orchestrates", NAVY)

    # Daemon loop hint
    _shape(slide, MSO_SHAPE.OVAL, Inches(0.2), Inches(2.35), Inches(9.6), Inches(4.55), LIGHT, TEAL)
    _label(slide, Inches(0.35), Inches(2.45), Inches(1.5), Inches(0.25), "pre_voyage_weather workflow", size=8, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

    # Row 2: ingest branch
    _arrow_v(slide, Inches(4.62), Inches(2.15), Inches(2.75), "ingest step")
    _agent(slide, Inches(3.85), Inches(2.75), "PreVoyage\nIngestAgent", "parse & plan", ORANGE)
    _arrow_h(slide, Inches(5.45), Inches(3.2), Inches(5.85), "upsert")
    _store(slide, Inches(5.85), Inches(2.95), "Voyage\nRegistry")
    _artifact(slide, Inches(7.35), Inches(2.55), "master", "route.json", SLATE)
    _arrow_h(slide, Inches(7.2), Inches(2.9), Inches(7.35))
    _label(slide, Inches(7.35), Inches(2.35), Inches(1.1), Inches(0.2), "voyage record", size=7, color=SLATE)

    _label(slide, Inches(5.75), Inches(3.95), Inches(1.55), Inches(0.55), "six_hour_plan\n(lat, lon, ETA)", size=7, color=PURPLE)

    # Row 3: weather branch
    _arrow_v(slide, Inches(4.62), Inches(3.7), Inches(4.55), "weather step")
    _agent(slide, Inches(3.85), Inches(4.55), "Weather\nAgent", "fetch & merge", BLUE)
    _arrow_h(slide, Inches(5.45), Inches(5.0), Inches(6.55))
    _api_cloud(slide, Inches(6.55), Inches(4.72))

    # Row 4: outputs
    _arrow_v(slide, Inches(4.62), Inches(5.5), Inches(5.95), "artifacts")
    _artifact(slide, Inches(2.0), Inches(5.95), "track", "weather.json", BLUE)
    _artifact(slide, Inches(3.35), Inches(5.95), "passage", "weather.txt", GREEN)
    _artifact(slide, Inches(4.7), Inches(5.95), "bad wx", ".json", RED)

    _arrow_h(slide, Inches(6.0), Inches(6.3), Inches(6.45))
    # Map viewer box
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.45), Inches(5.75), Inches(2.9), Inches(1.35), WHITE, TEAL)
    _shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(5.95), Inches(2.6), Inches(0.95), PALE, TEAL)
    # tiny route line on map
    for i, (lx, ly) in enumerate([(6.75, 6.55), (7.15, 6.35), (7.55, 6.5), (7.95, 6.25), (8.35, 6.45)]):
        _shape(slide, MSO_SHAPE.OVAL, Inches(lx), Inches(ly), Inches(0.08), Inches(0.08), RED if i == 2 else TEAL)
    _label(slide, Inches(6.45), Inches(6.72), Inches(2.9), Inches(0.3), "Route Viewer — weather on map", size=8, bold=True, color=NAVY)

    # Background agent
    _agent(slide, Inches(0.55), Inches(4.55), "WeatherReport\nAgent", "scheduled re-fetch", BLUE)
    _arrow_h(slide, Inches(2.15), Inches(5.0), Inches(3.85), "polls registry")
    _label(slide, Inches(0.4), Inches(5.55), Inches(2.0), Inches(0.35), "background agent", size=7, color=SLATE)


def slide2(prs):
    """Zoom: Excel → voyage → waypoints → weather data flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Information Flow", "What each agent reads, transforms, and writes")

    stages = [
        (Inches(0.4), GREEN, "① Excel", ".xlsx", "vessel · ports\nCP speed · waypoints"),
        (Inches(2.55), ORANGE, "② Ingest Agent", "parse", "creates voyage\nin registry"),
        (Inches(4.7), PURPLE, "③ Registry", "state", "voyage_number\nsix_hour_plan"),
        (Inches(6.85), BLUE, "④ Weather Agent", "fetch", "wind · wave · swell\nper waypoint"),
    ]
    for left, color, title, tag, detail in stages:
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.1), Inches(1.85), Inches(1.55), color, NAVY)
        _label(slide, left, Inches(1.2), Inches(1.85), Inches(0.3), title, size=10, bold=True, color=WHITE)
        _label(slide, left, Inches(1.5), Inches(1.85), Inches(0.22), tag, size=8, color=PALE)
        _label(slide, left + Inches(0.08), Inches(1.78), Inches(1.7), Inches(0.75), detail, size=7, color=WHITE)
        if left < Inches(6.85):
            _arrow_h(slide, left + Inches(1.85), Inches(1.85), left + Inches(2.15))

    # Middle: waypoint generation detail
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(3.0), Inches(9.2), Inches(2.15), LIGHT, TEAL)
    _label(slide, Inches(0.55), Inches(3.1), Inches(3.0), Inches(0.25), "Waypoint generation (inside Ingest Agent)", size=9, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # Route schematic
    pts = [(1.0, 4.55), (2.2, 4.2), (3.5, 4.6), (4.8, 4.15), (6.1, 4.5), (7.3, 4.25), (8.5, 4.55)]
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        _shape(slide, MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2 - x1), Inches(0.04), TEAL)
    for i, (x, y) in enumerate(pts):
        fill = ORANGE if i % 2 == 0 else TEAL
        _shape(slide, MSO_SHAPE.OVAL, Inches(x - 0.06), Inches(y - 0.06), Inches(0.12), Inches(0.12), fill, NAVY)
        if i % 2 == 0:
            _label(slide, Inches(x - 0.2), Inches(y + 0.12), Inches(0.5), Inches(0.2), f"+{i*6}h", size=6, color=SLATE)

    _label(slide, Inches(0.55), Inches(4.85), Inches(2.5), Inches(0.2), "master route", size=7, color=SLATE, align=PP_ALIGN.LEFT)
    _label(slide, Inches(7.8), Inches(4.85), Inches(1.5), Inches(0.2), "6h waypoints", size=7, color=TEAL, align=PP_ALIGN.RIGHT)
    _arrow_h(slide, Inches(3.2), Inches(3.55), Inches(5.5), "walk at CP speed → six_hour_plan[]")

    # Bottom: weather merge
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.45), Inches(4.2), Inches(1.55), WHITE, BLUE)
    _label(slide, Inches(0.55), Inches(5.55), Inches(1.5), Inches(0.25), "per waypoint", size=8, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    for i, (lbl, val, col) in enumerate([("wind", "28 kn", TEAL), ("wave", "3.2 m", BLUE), ("swell", "2.1 m", NAVY)]):
        lx = Inches(0.6 + i * 1.25)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, lx, Inches(5.9), Inches(1.05), Inches(0.75), col)
        _label(slide, lx, Inches(5.98), Inches(1.05), Inches(0.25), lbl, size=7, color=WHITE)
        _label(slide, lx, Inches(6.28), Inches(1.05), Inches(0.25), val, size=9, bold=True, color=WHITE)

    _arrow_h(slide, Inches(4.65), Inches(6.25), Inches(5.05), "merge")
    _artifact(slide, Inches(5.05), Inches(5.85), "voyage_track", "weather.json", BLUE)

    _arrow_h(slide, Inches(6.15), Inches(6.25), Inches(6.55), "")
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(5.45), Inches(3.05), Inches(1.55), WHITE, GREEN)
    _label(slide, Inches(6.7), Inches(5.55), Inches(1.2), Inches(0.25), "TXT report", size=8, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    _label(slide, Inches(6.7), Inches(5.9), Inches(2.7), Inches(0.9), "passage weather table\nhard weather flags\noperator summary", size=8, color=SLATE, align=PP_ALIGN.LEFT)


def slide3(prs):
    """Agentic model + continuous operation — visual only."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Why Agentic?", "Specialists · shared state · autonomous loops")

    # Center: registry hub
    cx, cy = _store(slide, Inches(4.1), Inches(2.55), "Voyage\nRegistry")

    agents = [
        (Inches(0.55), Inches(1.35), "InboxWatch\nAgent", TEAL, "watches\ninbox"),
        (Inches(0.55), Inches(3.85), "WeatherReport\nAgent", BLUE, "polls\ndue jobs"),
        (Inches(7.55), Inches(1.35), "PreVoyage\nIngestAgent", ORANGE, "creates\nvoyage"),
        (Inches(7.55), Inches(3.85), "Weather\nAgent", BLUE, "enriches\nplan"),
    ]
    for left, top, name, color, action in agents:
        _agent(slide, left, top, name.replace("\n", " "), action.replace("\n", " "), color)
        # arrow toward registry
        if left < Inches(4):
            _arrow_h(slide, left + Inches(1.55), top + Inches(0.48), Inches(4.1), "")
        else:
            _arrow_h(slide, Inches(5.45), top + Inches(0.48), left, "")

    # Daemon ring
    _shape(slide, MSO_SHAPE.OVAL, Inches(2.8), Inches(5.15), Inches(4.4), Inches(1.55), LIGHT, NAVY)
    _label(slide, Inches(2.8), Inches(5.55), Inches(4.4), Inches(0.35), "run_daemon.py — never stops, agents poll on timers", size=9, bold=True, color=NAVY)

    # Legend icons
    legend = [
        (Inches(0.45), "Agent", TEAL),
        (Inches(2.1), "Shared state", PURPLE),
        (Inches(3.95), "Artifact", SLATE),
        (Inches(5.55), "External API", BLUE),
    ]
    for left, text, col in legend:
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(6.95), Inches(0.35), Inches(0.22), col)
        _label(slide, left + Inches(0.42), Inches(6.9), Inches(1.2), Inches(0.3), text, size=8, color=SLATE, align=PP_ALIGN.LEFT)

    _label(
        slide,
        Inches(0.4),
        Inches(7.15),
        Inches(9.2),
        Inches(0.25),
        "Flow tag: pre_voyage_weather  →  ingest → weather  →  map + TXT report",
        size=9,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide1(prs)
    slide2(prs)
    slide3(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
